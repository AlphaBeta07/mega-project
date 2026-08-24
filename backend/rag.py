import os
import json
import csv
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
import io
import requests
from datetime import datetime
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import pandas as pd
from youtube_transcript_api import YouTubeTranscriptApi
import chromadb
from chromadb.utils import embedding_functions
from openai import AsyncOpenAI
from pypdf import PdfReader
import re
import asyncio
import edge_tts
import uuid
from pydub import AudioSegment
import urllib.parse
import aiohttp

# Initialize ChromaDB
chroma_client = chromadb.PersistentClient(path="./chroma_data")
emb_fn = embedding_functions.DefaultEmbeddingFunction()
collection = chroma_client.get_or_create_collection(name="studysnap_collection", embedding_function=emb_fn)

# Initialize LM Studio client
lm_studio_client = AsyncOpenAI(base_url="http://localhost:1234/v1", api_key="lm-studio")
MODEL_NAME = "my_custom_200m_model_gguf/AnishLandage/SnapStudyAI"

def extract_youtube_transcript(url: str) -> str:
    try:
        video_id = ""
        if "v=" in url:
            video_id = url.split("v=")[1].split("&")[0]
        elif "youtu.be/" in url:
            video_id = url.split("youtu.be/")[1].split("?")[0]
        elif "embed/" in url:
            video_id = url.split("embed/")[1].split("?")[0]
            
        if not video_id: 
            raise ValueError("Could not find video ID in URL.")
        
        ytt_api = YouTubeTranscriptApi()
        
        try:
            # Try getting standard English or auto-generated
            transcript_list = ytt_api.list(video_id)
            transcript = transcript_list.find_transcript(['en', 'en-US', 'en-GB', 'hi'])
            data = transcript.fetch()
        except:
            try:
                # Fallback to any auto-generated transcript if manual English isn't found
                transcript_list = ytt_api.list(video_id)
                transcript = list(transcript_list)[0] # Just get the first available one
                data = transcript.fetch()
            except Exception as e:
                # If all else fails, use the basic method
                data = ytt_api.fetch(video_id)
                
        return " ".join([t.text if hasattr(t, 'text') else t.get('text', '') for t in data])
    except Exception as e:
        print(f"YouTube Extraction Error: {e}")
        raise ValueError(f"Failed to process YouTube transcript: {str(e)}. The video might not have captions enabled.")

def extract_webpage(url: str) -> str:
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers)
        soup = BeautifulSoup(response.text, 'html.parser')
        # Remove scripts and styles
        for script in soup(["script", "style"]):
            script.extract()
        return soup.get_text(separator="\n", strip=True)
    except Exception as e:
        print(f"Webpage Extraction Error: {e}")
        return ""

import whisper

def extract_audio(file_path: str) -> str:
    """Extract text from audio using local Whisper instance."""
    try:
        # Load the base model ('tiny', 'base', 'small', 'medium', 'large')
        # 'tiny' is much faster for local transcription
        model = whisper.load_model("tiny")
        
        # Transcribe the audio file
        # Whisper automatically handles chunking for long audio files
        result = model.transcribe(file_path, fp16=False)
        
        return result.get("text", "")
    except Exception as e:
        print(f"Whisper Audio Extraction Error: {e}")
        return "Audio transcription failed. Please ensure ffmpeg is installed."

def extract_text(file_path: str, filename: str) -> str:
    """Extract text from all supported file types."""
    text = ""
    ext = filename.split('.')[-1].lower() if '.' in filename else ""
    
    try:
        if ext == 'pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                if page_text := page.extract_text(): text += page_text + "\n"
        elif ext in ['txt', 'md', 'html', 'xml']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                if ext in ['html', 'xml']:
                    soup = BeautifulSoup(content, 'html.parser')
                    text = soup.get_text(separator="\n", strip=True)
                else:
                    text = content
        elif ext == 'json':
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                text = json.dumps(data, indent=2)
        elif ext == 'csv':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    text += " ".join(row) + "\n"
        elif ext == 'xlsx':
            df = pd.read_excel(file_path)
            text = df.to_string()
        elif ext == 'docx':
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == 'pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ['wav', 'mp3', 'm4a', 'flac', 'ogg', 'opus', 'aac']:
            text = extract_audio(file_path)
    except Exception as e:
        print(f"Error extracting {filename}: {e}")
        
    return text

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200):
    chunks = []
    start = 0
    text_length = len(text)
    while start < text_length:
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
    return chunks

def ingest_document(file_path: str, filename: str, file_id: str):
    text = ""
    # Check if it's a URL or YouTube link based on filename format or type passed from main
    if filename.startswith("http"):
        if "youtube.com" in filename or "youtu.be" in filename:
            text = extract_youtube_transcript(filename)
            filename = f"YouTube: {filename}"
        else:
            text = extract_webpage(filename)
            filename = f"Web: {filename}"
    else:
        text = extract_text(file_path, filename)
        
    if not text.strip():
        raise ValueError("Could not extract any text from the document.")

    chunks = chunk_text(text)
    ids = [f"{file_id}_chunk_{i}" for i in range(len(chunks))]
    metadatas = [{"source": filename, "file_id": file_id} for _ in chunks]
    
    collection.add(documents=chunks, metadatas=metadatas, ids=ids)

    doc_meta = {
        "id": file_id,
        "filename": filename,
        "type": filename.split('.')[-1] if '.' in filename and not filename.startswith("http") else "url",
        "file_path": file_path,
        "raw_text": text[:5000],  # First 5000 chars preview
        "chunk_count": len(chunks)
    }

    # Persist in MongoDB Sync Database
    try:
        from database.connection import get_sync_db
        from database.config import settings
        sync_db = get_sync_db()
        sync_db[settings.SOURCES_COLLECTION].replace_one(
            {"id": file_id},
            {
                "id": file_id,
                "filename": filename,
                "type": doc_meta["type"],
                "file_path": file_path,
                "raw_text": text,
                "chunk_count": len(chunks),
                "created_at": datetime.utcnow()
            },
            upsert=True
        )
    except Exception as e:
        print(f"MongoDB Source Sync Warning: {e}")

    return doc_meta

def get_all_sources():
    # Try fetching from MongoDB first for full persistent data
    try:
        from database.connection import get_sync_db
        from database.config import settings
        sync_db = get_sync_db()
        cursor = sync_db[settings.SOURCES_COLLECTION].find({}, {"_id": 0, "raw_text": 0})
        sources = list(cursor)
        if sources:
            return sources
    except Exception as e:
        print(f"MongoDB Source Fetch Warning: {e}")

    # Fallback to ChromaDB
    results = collection.get(include=["metadatas"])
    metadatas = results.get("metadatas", [])
    unique_sources = {}
    for meta in metadatas:
        if meta["file_id"] not in unique_sources:
            unique_sources[meta["file_id"]] = {
                "id": meta["file_id"],
                "filename": meta["source"],
                "type": meta["source"].split('.')[-1] if '.' in meta["source"] and not meta["source"].startswith("http") else "url"
            }
    return list(unique_sources.values())

def delete_source(file_id: str):
    """Delete a source and all its chunks from ChromaDB and MongoDB."""
    collection.delete(where={"file_id": file_id})
    try:
        from database.connection import get_sync_db
        from database.config import settings
        sync_db = get_sync_db()
        sync_db[settings.SOURCES_COLLECTION].delete_one({"id": file_id})
    except Exception as e:
        print(f"MongoDB Source Delete Warning: {e}")

async def chat_with_context(query: str, history: list = None, selected_source_ids: list = None, response_language: str = "English"):
    if history is None: history = []
    
    where_clause = None
    if selected_source_ids is not None and len(selected_source_ids) > 0:
        if len(selected_source_ids) == 1:
            where_clause = {"file_id": selected_source_ids[0]}
        else:
            where_clause = {"file_id": {"$in": selected_source_ids}}
            
    # If a filter is provided, pass it to where
    if where_clause:
        results = collection.query(query_texts=[query], n_results=4, where=where_clause)
    else:
        results = collection.query(query_texts=[query], n_results=4)
        
    documents = results.get("documents", [[]])[0]
    metadatas = results.get("metadatas", [[]])[0]
    
    context = ""
    sources_used = []
    for doc, meta in zip(documents, metadatas):
        context += f"--- Source: {meta['source']} ---\n{doc}\n\n"
        if meta['source'] not in sources_used:
            sources_used.append(meta['source'])
            
    system_prompt = (
        f"You are StudySnap AI, a helpful and knowledgeable assistant. "
        f"Use the provided context to answer the user's question. "
        f"If the answer is not in the context, say that you don't know based on the provided documents.\n"
        f"IMPORTANT: You must respond in {response_language}.\n\n"
        f"Context:\n{context}"
    )
    
    messages = [{"role": "system", "content": system_prompt}]
    for msg in history:
        messages.append({"role": msg.get("role", "user"), "content": msg.get("content", "")})
        
    final_query = query
    if response_language and response_language.lower() != "english":
        lang_def = ""
        if response_language.lower() == "hinglish":
            lang_def = " (a mix of Hindi and English written in the English alphabet)"
        elif response_language.lower() == "manglish":
            lang_def = " (a mix of Malayalam and English written in the English alphabet)"
            
        final_query += f"\n\n[CRITICAL DIRECTIVE: You MUST write your ENTIRE response, including all headers, bullet points, and explanations, in {response_language}{lang_def}. Do NOT output standard English!]"

    messages.append({"role": "user", "content": final_query})
    
    response = await lm_studio_client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        temperature=0.7,
        max_tokens=1024,
    )
    
    answer = response.choices[0].message.content
    return answer, sources_used

async def generate_podcast_script(selected_source_ids: list = None, response_language: str = "English"):
    where_clause = None
    if selected_source_ids and len(selected_source_ids) > 0:
        if len(selected_source_ids) == 1:
            where_clause = {"file_id": selected_source_ids[0]}
        else:
            where_clause = {"file_id": {"$in": selected_source_ids}}
            
    # Get all chunks for these files, limit to 20 to avoid context limits
    if where_clause:
        results = collection.get(where=where_clause, limit=20)
    else:
        results = collection.get(limit=20)
        
    documents = results.get("documents", [])
    metadatas = results.get("metadatas", [])
    
    context = ""
    for doc, meta in zip(documents, metadatas):
        context += f"--- Source: {meta['source']} ---\n{doc}\n\n"
        
    lang_def = ""
    if response_language and response_language.lower() == "hinglish":
        lang_def = " (a mix of Hindi and English written in the English alphabet)"
    elif response_language and response_language.lower() == "manglish":
        lang_def = " (a mix of Malayalam and English written in the English alphabet)"

    system_prompt = (
        "You are an expert podcast producer. Create a multi-turn, engaging podcast script between two hosts, "
        "Host A and Host B, discussing the provided documents. "
        f"IMPORTANT: The hosts MUST speak and banter in {response_language}{lang_def}. Write the actual dialogue in this language.\n"
        "Host A should introduce the topic, and they should have a natural back-and-forth banter. "
        "Return the output STRICTLY as a JSON array of objects, with each object having 'speaker' ('Host A' or 'Host B') "
        "and 'text' (the spoken dialogue). Do not include any other text outside the JSON array.\n\n"
        "Context:\n" + context
    )
    
    messages = [{"role": "system", "content": system_prompt}]
    
    response = await lm_studio_client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        temperature=0.7,
        max_tokens=2048,
    )
    
    raw_response = response.choices[0].message.content
    
    try:
        # Extract JSON if wrapped in markdown code blocks
        json_match = re.search(r'\[.*\]', raw_response, re.DOTALL)
        if json_match:
            script_json = json.loads(json_match.group(0))
        else:
            script_json = json.loads(raw_response)
        return script_json
    except Exception as e:
        print(f"Error parsing script JSON: {e}")
        print(f"Raw Response: {raw_response}")
        return [{"speaker": "Host A", "text": "Welcome to our Audio Overview! We had a small issue generating the script."}, 
                {"speaker": "Host B", "text": "Yes, unfortunately we couldn't process the documents correctly."}]

async def generate_tts_audio(script_json):
    audio_files = []
    
    # Define voices
    voices = {
        "Host A": "en-US-ChristopherNeural",
        "Host B": "en-US-AriaNeural"
    }
    
    for i, line in enumerate(script_json):
        speaker = line.get("speaker", "Host A")
        text = line.get("text", "")
        voice = voices.get(speaker, voices["Host A"])
        
        output_file = f"temp_{uuid.uuid4()}_{i}.mp3"
        try:
            communicate = edge_tts.Communicate(text, voice)
            await communicate.save(output_file)
        except Exception as e:
            print(f"edge-tts failed: {e}. Falling back to gTTS.")
            from gtts import gTTS
            tts = gTTS(text=text, lang='en')
            tts.save(output_file)
        audio_files.append(output_file)
        
    # Merge audio files using pydub
    combined = AudioSegment.empty()
    for file in audio_files:
        try:
            segment = AudioSegment.from_mp3(file)
            combined += segment
            # Add a small pause between speakers (300ms)
            combined += AudioSegment.silent(duration=300)
        except Exception as e:
            print(f"Error merging {file}: {e}")
            
    final_output_name = f"podcast_{uuid.uuid4()}.mp3"
    final_output = f"uploads/{final_output_name}"
    combined.export(final_output, format="mp3")
    
    # Cleanup temp files
    for file in audio_files:
        if os.path.exists(file):
            os.remove(file)
            
    return final_output_name

async def generate_infographic_markdown(selected_source_ids: list, style: str, detail_level: str, custom_prompt: str, response_language: str = "English"):
    where_clause = None
    if selected_source_ids and len(selected_source_ids) > 0:
        if len(selected_source_ids) == 1:
            where_clause = {"file_id": selected_source_ids[0]}
        else:
            where_clause = {"file_id": {"$in": selected_source_ids}}
            
    if where_clause:
        results = collection.get(where=where_clause, limit=20)
    else:
        results = collection.get(limit=20)
        
    documents = results.get("documents", [])
    
    context = "\n\n".join(documents)
    
    lang_def = ""
    if response_language and response_language.lower() == "hinglish":
        lang_def = " (a mix of Hindi and English written in the English alphabet)"
    elif response_language and response_language.lower() == "manglish":
        lang_def = " (a mix of Malayalam and English written in the English alphabet)"

    system_prompt = f"""You are StudySnap AI's Infographic Engine, inspired by Google NotebookLM.

Analyze all uploaded sources and generate a professional educational infographic that is accurate, concise, and visually structured. Use only the provided content—never hallucinate or add unsupported information.

The infographic should include:
- Title & Overview
- Core Concepts
- Process Flow
- Architecture/Relationship Diagram (if applicable)
- Comparison Table (if applicable)
- Applications
- Formula/Important Facts
- Revision Box
- Key Takeaways

Automatically choose the best visual elements such as flowcharts, mind maps, timelines, architecture diagrams, comparison tables, icons, and illustrations.

Use a clean Google Material Design style:
- Minimal, professional layout
- White background with blue accents
- Flat vector illustrations
- Clear typography
- Well-spaced sections

Return the output in Markdown with clear headings.
The entire response MUST be in {response_language}{lang_def}.

At the end, generate detailed Gemini Image prompts for every illustration, diagram, and icon required to render the infographic.

User's custom instructions: {custom_prompt}

Context:
{context}"""
    
    try:
        model = genai.GenerativeModel('gemini-2.5-flash')
        response = model.generate_content(system_prompt)
        return response.text
    except Exception as e:
        print(f"Error generating infographic with Gemini: {e}")
        return f"# Error Generating Infographic\n\nThere was an error communicating with the Gemini API: {e}"

async def generate_mind_map_data(selected_source_ids: list, custom_prompt: str, response_language: str = "English"):
    where_clause = None
    if selected_source_ids and len(selected_source_ids) > 0:
        if len(selected_source_ids) == 1:
            where_clause = {"file_id": selected_source_ids[0]}
        else:
            where_clause = {"file_id": {"$in": selected_source_ids}}
            
    if where_clause:
        results = collection.get(where=where_clause, limit=20)
    else:
        results = collection.get(limit=20)
        
    documents = results.get("documents", [])
    context = "\n\n".join(documents)
    
    lang_def = ""
    if response_language and response_language.lower() == "hinglish":
        lang_def = " (a mix of Hindi and English written in the English alphabet)"
    elif response_language and response_language.lower() == "manglish":
        lang_def = " (a mix of Malayalam and English written in the English alphabet)"

    system_prompt = (
        "You are an expert knowledge architect. Your task is to analyze the provided documents and create a structured Mind Map. "
        f"User's custom instructions: {custom_prompt}\n"
        f"IMPORTANT: The node labels and context descriptions MUST be in {response_language}{lang_def}.\n\n"
        "Return the output STRICTLY as a JSON object with two keys: 'nodes' and 'edges'. "
        "Each node should have: "
        "- id: string (unique identifier) "
        "- data: object containing 'label' (short title) and 'context' (1-2 sentences explaining this node based on sources) "
        "- position: object containing 'x' and 'y' coordinates (arrange them hierarchically, e.g., root at 250,0, children spread out below). "
        "Each edge should have: "
        "- id: string (e.g., 'e1-2') "
        "- source: string (id of parent node) "
        "- target: string (id of child node). "
        "Do not include any markdown formatting, only pure JSON."
    )
    
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"Here is the source content:\n\n{context}"}
    ]
    
    response = await lm_studio_client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        temperature=0.3,
        max_tokens=2000,
    )
    
    try:
        import json
        content = response.choices[0].message.content.strip()
        if content.startswith("```json"):
            content = content[7:-3]
        elif content.startswith("```"):
            content = content[3:-3]
        return json.loads(content)
    except Exception as e:
        print("Failed to parse JSON mind map:", e)
        # Return a fallback mindmap
        return {
            "nodes": [
                {"id": "1", "data": {"label": "Error parsing map", "context": "Try again."}, "position": {"x": 250, "y": 250}}
            ],
            "edges": []
        }

import uuid
import edge_tts
from moviepy import AudioFileClip, ImageClip, concatenate_videoclips
from PIL import Image, ImageDraw, ImageFont

async def generate_video_overview(selected_source_ids: list = None, response_language: str = "English", custom_prompt: str = ""):
    where_clause = None
    if selected_source_ids and len(selected_source_ids) > 0:
        if len(selected_source_ids) == 1:
            where_clause = {"file_id": selected_source_ids[0]}
        else:
            where_clause = {"file_id": {"$in": selected_source_ids}}
            
    if where_clause:
        results = collection.get(where=where_clause, limit=20)
    else:
        results = collection.get(limit=20)
        
    documents = results.get("documents", [])
    
    context = ""
    for doc in documents:
        context += f"{doc}\n\n"
        
    system_prompt = f"""You are an educational video scriptwriter. 
Analyze the provided educational context.
Create a short, engaging video script explaining the core concepts.
Return ONLY a valid JSON array where each object represents a slide/scene.
The JSON array should have a maximum of 5 scenes to keep it short.
Do not wrap it in markdown blockquotes like ```json.
Format:
[
  {{
    "text_overlay": "Short title or key point for the screen",
    "narration": "The spoken explanation for this slide in {response_language}."
  }}
]

Context:
{context}

Custom Instructions: {custom_prompt}"""
    
    try:
        model = genai.GenerativeModel('gemini-2.5-flash')
        response = model.generate_content(system_prompt)
        text_resp = response.text.strip()
        if text_resp.startswith("```json"):
            text_resp = text_resp[7:]
        if text_resp.endswith("```"):
            text_resp = text_resp[:-3]
            
        scenes = json.loads(text_resp)
        
        video_id = str(uuid.uuid4())
        clips = []
        
        for i, scene in enumerate(scenes):
            narration = scene.get('narration', 'Continuing...')
            text_overlay = scene.get('text_overlay', '')
            
            # 1. Generate Audio
            audio_path = f"public/videos/temp_{video_id}_{i}.mp3"
            voice = "en-US-AriaNeural" if response_language.lower() == "english" else "es-ES-AlvaroNeural" # Defaulting for simplicity
            communicate = edge_tts.Communicate(narration, voice)
            await communicate.save(audio_path)
            
            # 2. Generate Image Slide (NotebookLM Style)
            img_path = f"public/videos/temp_{video_id}_{i}.png"
            # Background: Very dark blue/grey
            img = Image.new('RGB', (1920, 1080), color=(11, 12, 16))
            d = ImageDraw.Draw(img, 'RGBA')
            
            # Draw a subtle gradient orb/glow in the background
            orb_color = (66, 133, 244, 40) # Google Blue with low opacity
            d.ellipse([1920//2 - 500, 1080//2 - 500, 1920//2 + 500, 1080//2 + 500], fill=orb_color)
            
            # Draw Main Card
            card_margin_x = 250
            card_margin_y = 300
            card_rect = [card_margin_x, card_margin_y, 1920 - card_margin_x, 1080 - card_margin_y + 100]
            d.rounded_rectangle(card_rect, radius=40, fill=(26, 28, 35, 230), outline=(50, 50, 60, 255), width=2)
            
            # Draw Top Pill (Branding)
            pill_rect = [1920//2 - 150, 150, 1920//2 + 150, 210]
            d.rounded_rectangle(pill_rect, radius=30, fill=(40, 42, 50, 255), outline=(70, 70, 80, 255), width=2)
            
            try:
                font_large = ImageFont.truetype("arialbd.ttf", 64)
                font_small = ImageFont.truetype("arialbd.ttf", 24)
            except:
                font_large = ImageFont.load_default()
                font_small = ImageFont.load_default()
                
            # Pill text
            pill_text = "✨ StudySnap AI Overview"
            d.text((1920//2, 180), pill_text, fill=(200, 200, 210), font=font_small, anchor="mm")
                
            # Text wrapping logic for the card
            words = text_overlay.split()
            lines = []
            current_line = []
            for word in words:
                current_line.append(word)
                if len(" ".join(current_line)) > 35:
                    lines.append(" ".join(current_line))
                    current_line = []
            if current_line:
                lines.append(" ".join(current_line))
            wrapped_text = "\n".join(lines)
            
            # Draw Main Text centered in the card
            d.multiline_text((1920//2, (1080 + 100)//2), wrapped_text, fill=(240, 240, 245), font=font_large, anchor="mm", align="center", spacing=20)
            
            # Draw a subtle "Audio Playing" indicator at the bottom of the card
            indicator_text = "▶ Playing Audio Narration..."
            d.text((1920//2, 1080 - card_margin_y + 50), indicator_text, fill=(100, 100, 120), font=font_small, anchor="mm")
            
            # Apply blur to background orb (hacky way using ImageFilter on a base layer)
            try:
                from PIL import ImageFilter
                base = Image.new('RGB', (1920, 1080), color=(11, 12, 16))
                d_base = ImageDraw.Draw(base, 'RGBA')
                d_base.ellipse([1920//2 - 600, 1080//2 - 600, 1920//2 + 600, 1080//2 + 600], fill=(66, 133, 244, 40))
                d_base.ellipse([200, 200, 800, 800], fill=(168, 85, 247, 30)) # Purple orb
                base = base.filter(ImageFilter.GaussianBlur(150))
                base.paste(img, (0,0), img.convert('RGBA'))
                base.save(img_path)
            except:
                img.save(img_path)
            
            # 3. Combine into MoviePy Clip
            audio_clip = AudioFileClip(audio_path)
            img_clip = ImageClip(img_path).with_duration(audio_clip.duration)
            img_clip = img_clip.with_audio(audio_clip)
            clips.append(img_clip)
            
        # Concatenate all clips
        final_video = concatenate_videoclips(clips)
        output_path = f"public/videos/{video_id}.mp4"
        # 24 fps is enough for static slideshow
        final_video.write_videofile(output_path, fps=24, codec="libx264", audio_codec="aac")
        
        # Cleanup clips from memory
        final_video.close()
        for c in clips:
            c.close()
            
        # Clean temp files (optional, leaving them for debugging is fine or we can delete)
        for i in range(len(scenes)):
            if os.path.exists(f"public/videos/temp_{video_id}_{i}.mp3"):
                os.remove(f"public/videos/temp_{video_id}_{i}.mp3")
            if os.path.exists(f"public/videos/temp_{video_id}_{i}.png"):
                os.remove(f"public/videos/temp_{video_id}_{i}.png")
                
        return f"/videos/{video_id}.mp4"
        
    except Exception as e:
        print(f"Error generating video overview: {e}")
        return f"# Error Generating Video Overview\n\nThere was an error generating the video: {e}"
